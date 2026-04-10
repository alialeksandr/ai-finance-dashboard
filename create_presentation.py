#!/usr/bin/env python3
"""Generate CFO Dashboard consulting presentation."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
import os

# ── Colour palette (dark consulting style) ──────────────────────────
BG_DARK    = RGBColor(0x1B, 0x1F, 0x2B)   # slide background
BG_CARD    = RGBColor(0x23, 0x28, 0x38)   # card / box fill
ACCENT     = RGBColor(0x00, 0xC9, 0xA7)   # teal accent
ACCENT2    = RGBColor(0x6C, 0x63, 0xFF)   # purple accent
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xB0, 0xB8, 0xC8)
MID_GRAY   = RGBColor(0x7A, 0x82, 0x96)
ORANGE     = RGBColor(0xFF, 0x8C, 0x00)
RED_ACCENT = RGBColor(0xFF, 0x5C, 0x5C)
GREEN      = RGBColor(0x4C, 0xAF, 0x50)
YELLOW     = RGBColor(0xFF, 0xD5, 0x4F)
DARK_CARD2 = RGBColor(0x2A, 0x2F, 0x42)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width  = SLIDE_W
prs.slide_height = SLIDE_H

# ── Helper functions ─────────────────────────────────────────────────
def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_textbox(slide, left, top, width, height, text="", font_size=12,
                color=WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox

def add_rich_textbox(slide, left, top, width, height):
    """Return (textbox, text_frame) for manual paragraph building."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    return txBox, tf

def add_paragraph(tf, text, font_size=12, color=WHITE, bold=False,
                  alignment=PP_ALIGN.LEFT, space_after=Pt(4), level=0, font_name="Calibri"):
    if len(tf.paragraphs) == 1 and tf.paragraphs[0].text == "":
        p = tf.paragraphs[0]
    else:
        p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    p.space_after = space_after
    p.level = level
    return p

def add_rounded_rect(slide, left, top, width, height, fill_color=BG_CARD, line_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    # reduce corner radius
    shape.adjustments[0] = 0.05
    return shape

def add_line(slide, x1, y1, x2, y2, color=ACCENT, width=Pt(2)):
    connector = slide.shapes.add_connector(1, x1, y1, x2, y2)  # MSO_CONNECTOR.STRAIGHT
    connector.line.color.rgb = color
    connector.line.width = width
    return connector

def add_circle(slide, left, top, size, fill_color):
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, size, size)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape

def add_milestone_diamond(slide, cx, cy, size, fill_color):
    shape = slide.shapes.add_shape(MSO_SHAPE.DIAMOND, cx - size//2, cy - size//2, size, size)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape

# ── Footer helper ────────────────────────────────────────────────────
def add_footer(slide, page_num, total=4):
    add_textbox(slide, Inches(0.5), Inches(7.05), Inches(4), Inches(0.35),
                "CFO Dashboard  |  Confidential", font_size=8, color=MID_GRAY)
    add_textbox(slide, Inches(11), Inches(7.05), Inches(2), Inches(0.35),
                f"{page_num} / {total}", font_size=8, color=MID_GRAY, alignment=PP_ALIGN.RIGHT)

# ======================================================================
# SLIDE 1 – CFO Dashboard – Vision
# ======================================================================
slide1 = prs.slides.add_slide(prs.slide_layouts[6])  # blank
set_slide_bg(slide1, BG_DARK)

# Top accent line
add_line(slide1, Inches(0), Inches(0.05), SLIDE_W, Inches(0.05), ACCENT, Pt(4))

# Title
add_textbox(slide1, Inches(0.7), Inches(0.3), Inches(12), Inches(0.6),
            "CFO Dashboard \u2013 Vision", font_size=32, color=WHITE, bold=True)

# Subtitle line
add_textbox(slide1, Inches(0.7), Inches(0.85), Inches(12), Inches(0.35),
            "AI-enabled dashboarding & BI for financial professionals", font_size=14, color=ACCENT)

# Divider
add_line(slide1, Inches(0.7), Inches(1.25), Inches(12.6), Inches(1.25), MID_GRAY, Pt(0.5))

# ── LEFT COLUMN: Vision ──────────────────────────────────────────────
left_card = add_rounded_rect(slide1, Inches(0.7), Inches(1.45), Inches(5.8), Inches(5.3), BG_CARD)

# "VISION" label
add_textbox(slide1, Inches(1.0), Inches(1.55), Inches(2), Inches(0.35),
            "VISION", font_size=11, color=ACCENT, bold=True)

# Date
add_textbox(slide1, Inches(4.3), Inches(1.55), Inches(2), Inches(0.35),
            "Discussed 27/03", font_size=10, color=MID_GRAY, alignment=PP_ALIGN.RIGHT)

# Vision body
_, tf = add_rich_textbox(slide1, Inches(1.0), Inches(1.95), Inches(5.2), Inches(1.8))
add_paragraph(tf, "Discover possibilities to create an AI-enabled dashboarding and BI solution for:",
              font_size=12, color=LIGHT_GRAY, space_after=Pt(8))
add_paragraph(tf, "\u2022  Financial professionals working with SMBs", font_size=11, color=WHITE, space_after=Pt(4), level=1)
add_paragraph(tf, "\u2022  SMB owners", font_size=11, color=WHITE, space_after=Pt(4), level=1)
add_paragraph(tf, "\u2022  Second stage: payroll companies", font_size=11, color=WHITE, space_after=Pt(8), level=1)
add_paragraph(tf, "Global goal: commercialise as SaaS / micro-SaaS", font_size=12, color=ACCENT, bold=True, space_after=Pt(6))

# Key features box
feat_card = add_rounded_rect(slide1, Inches(1.0), Inches(4.0), Inches(5.2), Inches(2.5), DARK_CARD2, ACCENT)

add_textbox(slide1, Inches(1.2), Inches(4.1), Inches(3), Inches(0.3),
            "KEY FEATURES", font_size=10, color=ACCENT, bold=True)

_, tff = add_rich_textbox(slide1, Inches(1.2), Inches(4.45), Inches(4.8), Inches(2.0))
add_paragraph(tff, "\u2714  AI analytics \u2013 ask questions & generate management action insights",
              font_size=11, color=WHITE, space_after=Pt(6))
add_paragraph(tff, "\u2714  Advanced data slicing / dicing",
              font_size=11, color=WHITE, space_after=Pt(6))
add_paragraph(tff, "\u2714  Traceability & audit trail",
              font_size=11, color=WHITE, space_after=Pt(6))

# ── RIGHT COLUMN: References ─────────────────────────────────────────
right_card = add_rounded_rect(slide1, Inches(6.8), Inches(1.45), Inches(6.0), Inches(5.3), BG_CARD)

add_textbox(slide1, Inches(7.1), Inches(1.55), Inches(3), Inches(0.35),
            "REFERENCES", font_size=11, color=ACCENT2, bold=True)

# Place reference images (2x2 grid)
img_dir = "/home/user/ai-finance-dashboard"
ref_images = []
for f in sorted(os.listdir("/tmp")):
    if f.endswith(('.png', '.jpg', '.jpeg', '.webp')):
        ref_images.append(os.path.join("/tmp", f))

# Check uploads directory too
for d in ["/home/user/uploads", "/home/user"]:
    if os.path.isdir(d):
        for f in sorted(os.listdir(d)):
            if f.endswith(('.png', '.jpg', '.jpeg', '.webp')) and os.path.join(d, f) not in ref_images:
                ref_images.append(os.path.join(d, f))

# We'll place placeholder boxes for reference screenshots
img_positions = [
    (Inches(7.1), Inches(2.0), Inches(2.7), Inches(1.8)),
    (Inches(10.0), Inches(2.0), Inches(2.7), Inches(1.8)),
    (Inches(7.1), Inches(4.0), Inches(2.7), Inches(1.8)),
    (Inches(10.0), Inches(4.0), Inches(2.7), Inches(1.8)),
]

ref_labels = ["AI Strategy Lab", "Financial Controller", "Executive Dashboard", "Data Upload"]

for i, (l, t, w, h) in enumerate(img_positions):
    box = add_rounded_rect(slide1, l, t, w, h, DARK_CARD2, MID_GRAY)
    add_textbox(slide1, l + Inches(0.15), t + Inches(0.6), w - Inches(0.3), Inches(0.4),
                f"[{ref_labels[i]}]", font_size=10, color=MID_GRAY, alignment=PP_ALIGN.CENTER)
    add_textbox(slide1, l + Inches(0.15), t + Inches(1.0), w - Inches(0.3), Inches(0.4),
                "Reference screenshot", font_size=8, color=MID_GRAY, alignment=PP_ALIGN.CENTER)

# Note about competitive landscape
_, tfn = add_rich_textbox(slide1, Inches(7.1), Inches(5.95), Inches(5.5), Inches(0.65))
add_paragraph(tfn, "\u26A0  Observation: significant number of operational solutions with similar functionality already exist in the market.",
              font_size=10, color=YELLOW, space_after=Pt(0))

add_footer(slide1, 1)


# ======================================================================
# SLIDE 2 – Roadmap and PoC
# ======================================================================
slide2 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide2, BG_DARK)
add_line(slide2, Inches(0), Inches(0.05), SLIDE_W, Inches(0.05), ACCENT, Pt(4))

add_textbox(slide2, Inches(0.7), Inches(0.3), Inches(12), Inches(0.6),
            "Roadmap & Proof of Concept", font_size=32, color=WHITE, bold=True)
add_line(slide2, Inches(0.7), Inches(0.95), Inches(12.6), Inches(0.95), MID_GRAY, Pt(0.5))

# ── ROADMAP (upper half) ─────────────────────────────────────────────
roadmap_card = add_rounded_rect(slide2, Inches(0.7), Inches(1.15), Inches(12.0), Inches(2.6), BG_CARD)

add_textbox(slide2, Inches(1.0), Inches(1.25), Inches(3), Inches(0.3),
            "ROADMAP", font_size=11, color=ACCENT, bold=True)

# Timeline bar
bar_y = Inches(2.15)
bar_left = Inches(1.2)
bar_right = Inches(12.0)
bar_width = bar_right - bar_left

# Background bar
bar_bg = slide2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, bar_left, bar_y, bar_width, Inches(0.15))
bar_bg.fill.solid()
bar_bg.fill.fore_color.rgb = DARK_CARD2
bar_bg.line.fill.background()
bar_bg.adjustments[0] = 0.5

# Progress fill (up to PoC ~ 40%)
bar_fill = slide2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, bar_left, bar_y, Emu(int(bar_width * 0.25)), Inches(0.15))
bar_fill.fill.solid()
bar_fill.fill.fore_color.rgb = ACCENT
bar_fill.line.fill.background()
bar_fill.adjustments[0] = 0.5

# Milestone markers and labels
milestones = [
    (0.0,  "Start",               "30.03", ACCENT),
    (0.15, "PoC Discovery",       "30.03 \u2013 04.10", LIGHT_GRAY),
    (0.30, "PoC Milestone",       "04.10", ACCENT),
    (0.50, "MVP Development",     "04.10 \u2013 31.05", LIGHT_GRAY),
    (0.42, "AI Fin Club Demo",    "01.05", ORANGE),
    (0.70, "MVP Milestone",       "31.05", ACCENT2),
    (0.82, "Scaling",             "31.05 \u2013 ...", LIGHT_GRAY),
    (1.0,  "First Prod Deploy",   "31.08", GREEN),
]

bar_left_val = 1.2
bar_width_inches = 10.8

for pct, label, date, clr in milestones:
    x_pos = Inches(bar_left_val + pct * bar_width_inches)

    # Diamond marker
    d_size = Inches(0.2)
    add_milestone_diamond(slide2, x_pos, bar_y + Inches(0.075), d_size, clr)

    # Label above
    add_textbox(slide2, x_pos - Inches(0.6), bar_y - Inches(0.55), Inches(1.5), Inches(0.35),
                label, font_size=8, color=clr, bold=True, alignment=PP_ALIGN.CENTER)
    # Date below
    add_textbox(slide2, x_pos - Inches(0.6), bar_y + Inches(0.3), Inches(1.5), Inches(0.25),
                date, font_size=7, color=MID_GRAY, alignment=PP_ALIGN.CENTER)

# ── PoC SCOPE (lower half) ───────────────────────────────────────────
scope_card = add_rounded_rect(slide2, Inches(0.7), Inches(3.95), Inches(12.0), Inches(1.6), BG_CARD, YELLOW)

add_textbox(slide2, Inches(1.0), Inches(4.0), Inches(3), Inches(0.3),
            "POC SCOPE RATIONALE", font_size=11, color=YELLOW, bold=True)

_, tfs = add_rich_textbox(slide2, Inches(1.0), Inches(4.35), Inches(11.5), Inches(1.1))
add_paragraph(tfs, "Original intention was to build something as close to references as possible, but with limited time and unavailable real-life data for hypothesis testing, the PoC was minimised to a smaller scope:",
              font_size=10, color=LIGHT_GRAY, space_after=Pt(6))
add_paragraph(tfs, "\u2192  Address one small feature fixing common SMB finance workflow issues: scattered data consolidation, repeating reconciliation tasks with unstructured data, lack of quick visualisation in small automation tools.",
              font_size=10, color=WHITE, bold=True, space_after=Pt(0))

# ── Goals achieved ───────────────────────────────────────────────────
add_textbox(slide2, Inches(0.7), Inches(5.7), Inches(12), Inches(0.35),
            "POC GOALS ACHIEVED", font_size=11, color=ACCENT, bold=True)

# Two goal cards side by side
alex_card = add_rounded_rect(slide2, Inches(0.7), Inches(6.05), Inches(5.8), Inches(1.2), BG_CARD, ACCENT)
add_textbox(slide2, Inches(1.0), Inches(6.1), Inches(2), Inches(0.3),
            "FOR ALEX", font_size=10, color=ACCENT, bold=True)
_, tfa = add_rich_textbox(slide2, Inches(1.0), Inches(6.4), Inches(5.3), Inches(0.8))
add_paragraph(tfa, "\u2022  Deep dive into problem domain & industry awareness", font_size=10, color=WHITE, space_after=Pt(3))
add_paragraph(tfa, "\u2022  Awareness of industry software landscape", font_size=10, color=WHITE, space_after=Pt(3))
add_paragraph(tfa, "\u2022  Applied new AI tools & known tools in new context", font_size=10, color=WHITE, space_after=Pt(0))

manny_card = add_rounded_rect(slide2, Inches(6.9), Inches(6.05), Inches(5.8), Inches(1.2), BG_CARD, ACCENT2)
add_textbox(slide2, Inches(7.2), Inches(6.1), Inches(2), Inches(0.3),
            "FOR MANNY", font_size=10, color=ACCENT2, bold=True)
_, tfm = add_rich_textbox(slide2, Inches(7.2), Inches(6.4), Inches(5.3), Inches(0.8))
add_paragraph(tfm, "\u2022  Effort estimation for potential full-scope implementation", font_size=10, color=WHITE, space_after=Pt(3))
add_paragraph(tfm, "\u2022  High-level architecture & technical foundation", font_size=10, color=WHITE, space_after=Pt(3))
add_paragraph(tfm, "\u2022  Potential material for AI Fin Club presentation", font_size=10, color=WHITE, space_after=Pt(0))

add_footer(slide2, 2)


# ======================================================================
# SLIDE 3 – Use Case & Architecture
# ======================================================================
slide3 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide3, BG_DARK)
add_line(slide3, Inches(0), Inches(0.05), SLIDE_W, Inches(0.05), ACCENT, Pt(4))

add_textbox(slide3, Inches(0.7), Inches(0.3), Inches(12), Inches(0.6),
            "Use Case & Architecture", font_size=32, color=WHITE, bold=True)
add_line(slide3, Inches(0.7), Inches(0.95), Inches(12.6), Inches(0.95), MID_GRAY, Pt(0.5))

# ── USE CASE DIAGRAM (upper half) ────────────────────────────────────
uc_card = add_rounded_rect(slide3, Inches(0.7), Inches(1.15), Inches(12.0), Inches(2.8), BG_CARD)
add_textbox(slide3, Inches(1.0), Inches(1.25), Inches(3), Inches(0.3),
            "USE CASE FLOW", font_size=11, color=ACCENT, bold=True)

# Flow diagram using shapes
flow_items = [
    ("User triggers\nDashboard Update", ACCENT2),
    ("Sales / Inventory\nReport from Gmail", RGBColor(0xEA, 0x43, 0x35)),
    ("Invoice Data\nfrom Xero", RGBColor(0x13, 0xB5, 0xEA)),
    ("Data Combined\n& Verified", ACCENT),
    ("Dashboard:\nOutstanding\nCollectables", GREEN),
]

x_start = Inches(1.2)
y_center = Inches(2.5)
box_w = Inches(1.9)
box_h = Inches(1.0)
gap = Inches(0.5)

for i, (label, clr) in enumerate(flow_items):
    x = x_start + i * (box_w + gap)
    box = add_rounded_rect(slide3, x, y_center - box_h//2, box_w, box_h, DARK_CARD2, clr)

    tb = add_textbox(slide3, x + Inches(0.1), y_center - box_h//2 + Inches(0.15),
                     box_w - Inches(0.2), box_h - Inches(0.2),
                     label, font_size=9, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

    # Arrow between boxes (except last)
    if i < len(flow_items) - 1:
        arrow_x = x + box_w
        add_textbox(slide3, arrow_x + Inches(0.05), y_center - Inches(0.15),
                    Inches(0.4), Inches(0.3), "\u25B6", font_size=16, color=clr, alignment=PP_ALIGN.CENTER)

# ── ARCHITECTURE DIAGRAM (lower half) ────────────────────────────────
arch_card = add_rounded_rect(slide3, Inches(0.7), Inches(4.15), Inches(12.0), Inches(3.0), BG_CARD)
add_textbox(slide3, Inches(1.0), Inches(4.25), Inches(3), Inches(0.3),
            "ARCHITECTURE", font_size=11, color=ACCENT, bold=True)

# Architecture components
arch_components = [
    # (x, y, w, h, label, color)
    (Inches(1.2), Inches(5.1), Inches(1.5), Inches(0.8), "Gmail\n(Data Source)", RGBColor(0xEA, 0x43, 0x35)),
    (Inches(1.2), Inches(6.1), Inches(1.5), Inches(0.8), "Xero\n(Accounting)", RGBColor(0x13, 0xB5, 0xEA)),
    (Inches(3.6), Inches(5.4), Inches(1.8), Inches(1.2), "n8n\n(Workflow\nOrchestration)", ACCENT2),
    (Inches(6.3), Inches(4.8), Inches(2.2), Inches(0.8), "Front-end\nJS, Vue, npm, nginx", ACCENT),
    (Inches(6.3), Inches(5.9), Inches(2.2), Inches(0.8), "Back-end\nJS, npm", ACCENT),
    (Inches(9.3), Inches(5.1), Inches(2.5), Inches(1.2), "AI SMB CFO\nDashboard\n(ai-smb-cfo.nl)", GREEN),
]

for x, y, w, h, label, clr in arch_components:
    add_rounded_rect(slide3, x, y, w, h, DARK_CARD2, clr)
    add_textbox(slide3, x + Inches(0.1), y + Inches(0.08), w - Inches(0.2), h - Inches(0.15),
                label, font_size=9, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

# Arrows between arch components
arrow_pairs = [
    (Inches(2.7), Inches(5.5), Inches(3.6), Inches(5.7)),  # Gmail -> n8n
    (Inches(2.7), Inches(6.5), Inches(3.6), Inches(6.2)),  # Xero -> n8n
    (Inches(5.4), Inches(5.7), Inches(6.3), Inches(5.2)),  # n8n -> Frontend
    (Inches(5.4), Inches(6.0), Inches(6.3), Inches(6.3)),  # n8n -> Backend
    (Inches(8.5), Inches(5.2), Inches(9.3), Inches(5.5)),  # Frontend -> Dashboard
    (Inches(8.5), Inches(6.3), Inches(9.3), Inches(5.9)),  # Backend -> Dashboard
]
for x1, y1, x2, y2 in arrow_pairs:
    add_line(slide3, x1, y1, x2, y2, MID_GRAY, Pt(1.5))

# WEB-app label
add_textbox(slide3, Inches(6.3), Inches(6.8), Inches(2.2), Inches(0.3),
            "WEB-app", font_size=9, color=MID_GRAY, alignment=PP_ALIGN.CENTER)

# Note about architecture picture
add_textbox(slide3, Inches(1.0), Inches(7.0), Inches(5), Inches(0.3),
            "[Refer to attached architecture diagram for n8n workflow detail]", font_size=8, color=MID_GRAY)

add_footer(slide3, 3)


# ======================================================================
# SLIDE 4 – Open Questions & Other Business
# ======================================================================
slide4 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide4, BG_DARK)
add_line(slide4, Inches(0), Inches(0.05), SLIDE_W, Inches(0.05), ACCENT, Pt(4))

add_textbox(slide4, Inches(0.7), Inches(0.3), Inches(12), Inches(0.6),
            "Open Questions & Other Business", font_size=32, color=WHITE, bold=True)
add_line(slide4, Inches(0.7), Inches(0.95), Inches(12.6), Inches(0.95), MID_GRAY, Pt(0.5))

# Questions section
q_items = [
    {
        "q": "What are the real value features to build?",
        "detail": "This impacts functional architecture, integrations, and technical stack.",
        "icon": "\u2753",
        "color": ACCENT,
    },
    {
        "q": "Share contact of Payroll owner",
        "detail": "Check potential connection and explore payroll use case as second stage.",
        "icon": "\U0001F465",
        "color": ACCENT2,
    },
    {
        "q": "Alternative: develop plug-in for existing platform (e.g. Xero)?",
        "detail": "Instead of standalone SaaS, consider integrating into established ecosystems.",
        "icon": "\U0001F50C",
        "color": ORANGE,
    },
    {
        "q": "Reminder: Supplements topic",
        "detail": "Follow up on supplements discussion from previous meetings.",
        "icon": "\U0001F4CB",
        "color": YELLOW,
    },
]

y_pos = Inches(1.3)
for item in q_items:
    card = add_rounded_rect(slide4, Inches(0.7), y_pos, Inches(12.0), Inches(1.2), BG_CARD, item["color"])

    # Icon circle
    add_circle(slide4, Inches(1.1), y_pos + Inches(0.25), Inches(0.6), DARK_CARD2)
    add_textbox(slide4, Inches(1.1), y_pos + Inches(0.3), Inches(0.6), Inches(0.5),
                item["icon"], font_size=20, color=item["color"], alignment=PP_ALIGN.CENTER)

    # Question text
    add_textbox(slide4, Inches(2.0), y_pos + Inches(0.2), Inches(10.2), Inches(0.4),
                item["q"], font_size=16, color=WHITE, bold=True)
    add_textbox(slide4, Inches(2.0), y_pos + Inches(0.65), Inches(10.2), Inches(0.4),
                item["detail"], font_size=11, color=LIGHT_GRAY)

    y_pos += Inches(1.4)

# "Next steps" prompt
add_rounded_rect(slide4, Inches(3.5), Inches(6.5), Inches(6.3), Inches(0.65), DARK_CARD2, ACCENT)
add_textbox(slide4, Inches(3.5), Inches(6.55), Inches(6.3), Inches(0.5),
            "Discussion needed to align on priorities and next steps",
            font_size=14, color=ACCENT, bold=True, alignment=PP_ALIGN.CENTER)

add_footer(slide4, 4)

# ── SAVE ─────────────────────────────────────────────────────────────
output_path = "/home/user/ai-finance-dashboard/CFO_Dashboard_Presentation.pptx"
prs.save(output_path)
print(f"Presentation saved to: {output_path}")
